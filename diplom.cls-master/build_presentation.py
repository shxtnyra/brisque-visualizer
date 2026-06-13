# -*- coding: utf-8 -*-
"""
Генерация презентации к защите диплома (12 слайдов, 16:9).
Запуск:  python build_presentation.py
Результат: presentation-defense.pptx

Формулы рендерятся через matplotlib (mathtext) в PNG и вставляются как картинки.
Спецсимволы заданы через \\u-escape, чтобы файл читался в любой кодировке.
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- спецсимволы ---
MU   = "\u03bc"
SIG  = "\u03c3"
SIGH = "\u03c3\u0302"
NU   = "\u03bd"
LAM  = "\u03bb"
ALP  = "\u03b1"
ETA  = "\u03b7"
GAM  = "\u0393"
RAD  = "\u221a"
TIMES = "\u00d7"
SUB1 = "\u2081"
SUB2 = "\u2082"
SUP2 = "\u00b2"
ARR  = "\u2192"
GG   = "\u226b"

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
FORM_DIR = os.path.join(HERE, "_formulas")
os.makedirs(FORM_DIR, exist_ok=True)
OUT = os.path.join(HERE, "presentation-defense.pptx")

# ---- Палитра -------------------------------------------------------------
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
INK    = RGBColor(0x22, 0x2A, 0x33)
MUTE   = RGBColor(0x5A, 0x66, 0x72)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD5, 0xDD, 0xE5)

FONT = "Calibri"
EMUW, EMUH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMUW
prs.slide_height = EMUH
BLANK = prs.slide_layouts[6]


# ---- Рендер формул --------------------------------------------------------
def render_formula(key, tex, fontsize=30, color="#1E3A5F"):
    path = os.path.join(FORM_DIR, key + ".png")
    fig = plt.figure(figsize=(0.1, 0.1))
    fig.patch.set_alpha(0)
    fig.text(0.0, 0.0, tex, fontsize=fontsize, color=color)
    fig.savefig(path, dpi=300, transparent=True, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    w, h = Image.open(path).size
    return path, w, h


def place_formula(s, info, rx, ry, rw, rh, max_h_in, max_w_frac=0.95):
    path, wpx, hpx = info
    max_w = int(rw * max_w_frac)
    iw = int(Inches(max_h_in) * wpx / hpx)
    ih = int(Inches(max_h_in))
    if iw > max_w:
        iw = max_w
        ih = int(iw * hpx / wpx)
    ix = int(rx + (rw - iw) / 2)
    iy = int(ry + (rh - ih) / 2)
    s.shapes.add_picture(path, ix, iy, width=iw, height=ih)


# ---- Базовые помощники ----------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(s, x, y, w, h, color, shape_type=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape_type, x, y, w, h)
    fill(sp, color)
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    return tb, tf


def set_run(r, text, size, color=INK, bold=False, italic=False, font=FONT):
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color


def para(tf, text, size, color=INK, bold=False, italic=False, space_after=6,
         level=0, align=PP_ALIGN.LEFT, bullet=False, font=FONT, line_spacing=1.05):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    p.line_spacing = line_spacing
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    _bullet(p, bullet)
    return p


def _bullet(p, on):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    if on:
        pPr.set("marL", str(Emu(Inches(0.32))))
        pPr.set("indent", str(-Emu(Inches(0.32))))
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2014"})
        pPr.append(buFont)
        pPr.append(buChar)
    else:
        buNone = pPr.makeelement(qn("a:buNone"), {})
        pPr.append(buNone)


def header(s, title, kicker=None, number=None):
    add_rect(s, 0, 0, EMUW, Inches(1.12), NAVY)
    add_rect(s, 0, Inches(1.12), EMUW, Inches(0.06), ACCENT)
    tb, tf = textbox(s, Inches(0.55), Inches(0.12), Inches(11.5), Inches(0.95),
                     anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        para(tf, kicker, 12, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True, space_after=2)
    para(tf, title, 27, color=WHITE, bold=True, space_after=0, line_spacing=1.0)
    if number is not None:
        nb, ntf = textbox(s, Inches(12.4), Inches(0.12), Inches(0.7), Inches(0.95),
                          anchor=MSO_ANCHOR.MIDDLE)
        para(ntf, number, 13, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True,
             align=PP_ALIGN.RIGHT, space_after=0)


def footer(s):
    tb, tf = textbox(s, Inches(0.55), Inches(7.04), Inches(12.2), Inches(0.36))
    para(tf, "Кунашев А. С.  ·  ИВТ-41БО  ·  Визуализация алгоритмов оценки качества изображений",
         9, color=MUTE, space_after=0)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ===========================================================================
# Слайд 1 — Титул
# ===========================================================================
s = slide()
add_rect(s, 0, 0, EMUW, EMUH, NAVY)
add_rect(s, 0, 0, Inches(0.22), EMUH, ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.5))
para(tf, "Ярославский государственный университет им. П. Г. Демидова",
     14, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True, space_after=0)
para(tf, "Факультет ИВТ  ·  Выпускная квалификационная работа",
     12, color=RGBColor(0x8F, 0xAB, 0xC9), space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(2.6),
                 anchor=MSO_ANCHOR.TOP)
para(tf, "Разработка инструмента для визуализации\nи прототипирования алгоритмов\nоценки качества изображений",
     34, color=WHITE, bold=True, space_after=0, line_spacing=1.08)
add_rect(s, Inches(0.92), Inches(5.05), Inches(2.6), Inches(0.05), ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(5.35), Inches(11), Inches(1.6))
para(tf, "Студент:  А. С. Кунашев, группа ИВТ-41БО", 16, color=WHITE, space_after=5)
para(tf, "Научный руководитель:  ст. преп. А. Г. Седов", 16,
     color=RGBColor(0xCF, 0xDD, 0xEC), space_after=5)
para(tf, "Ярославль, 2026", 14, color=RGBColor(0x8F, 0xAB, 0xC9), space_after=0)
notes(s, "Добрый день. Тема моей работы — инструмент для визуализации алгоритмов "
         "оценки качества изображений. Доклад займёт около семи минут.")

# ===========================================================================
# Слайд 2 — Зачем нужен инструмент (+ задачи)
# ===========================================================================
s = slide()
header(s, "Зачем нужен инструмент", kicker="ПРОБЛЕМА, ЦЕЛЬ И ЗАДАЧИ", number="2")
col_h = Inches(3.45)
add_rect(s, Inches(0.55), Inches(1.45), Inches(5.95), col_h, LIGHT)
add_rect(s, Inches(0.55), Inches(1.45), Inches(0.1), col_h, MUTE)
tb, tf = textbox(s, Inches(0.85), Inches(1.62), Inches(5.5), Inches(3.1))
para(tf, "ПРОБЛЕМА", 13, color=MUTE, bold=True, space_after=8)
para(tf, "Метрика IQA — цепочка осмысленных этапов: нормализация, признаки, регрессия",
     15, bullet=True, space_after=8)
para(tf, "Но реализации в OpenCV / MATLAB возвращают лишь итоговое число",
     15, bullet=True, space_after=8)
para(tf, "Промежуточные карты, распределения и параметры моделей скрыты — алгоритм трудно понять",
     15, bullet=True, space_after=0)
add_rect(s, Inches(6.8), Inches(1.45), Inches(5.95), col_h, NAVY)
add_rect(s, Inches(6.8), Inches(1.45), Inches(0.1), col_h, ACCENT)
tb, tf = textbox(s, Inches(7.1), Inches(1.62), Inches(5.5), Inches(3.1))
para(tf, "ЦЕЛЬ", 13, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True, space_after=8)
para(tf, "Инструмент для визуализации и изучения алгоритмов IQA",
     16, color=WHITE, bullet=True, space_after=8)
para(tf, "Снимок " + ARR + " ROI " + ARR + " поэтапный расчёт: карты, гистограммы, признаки, score",
     15, color=WHITE, bullet=True, space_after=8)
para(tf, "Модульная архитектура — метрики подключаются как модули",
     15, color=WHITE, bullet=True, space_after=8)
para(tf, "BRISQUE — первый демонстрационный модуль (пример, а не цель работы)",
     13, color=RGBColor(0xCF, 0xDD, 0xEC), bullet=True, space_after=0, level=1)
# задачи (бывший отдельный слайд) — одной лентой
add_rect(s, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.5), RGBColor(0xE8, 0xF1, 0xFA))
add_rect(s, Inches(0.55), Inches(5.2), Inches(0.1), Inches(1.5), ACCENT)
tb, tf = textbox(s, Inches(0.85), Inches(5.32), Inches(11.8), Inches(1.3))
para(tf, "ЗАДАЧИ РАБОТЫ", 12, color=NAVY, bold=True, space_after=6)
chain_tasks = ("Обзор IQA  " + ARR + "  выбор этапов BRISQUE  " + ARR +
               "  стек и архитектура  " + ARR + "  ядро в Web Worker  " + ARR +
               "  UI: карты, графики, признаки  " + ARR + "  тесты с MATLAB")
para(tf, chain_tasks, 15, color=INK, bold=True, space_after=0, line_spacing=1.1)
footer(s)
notes(s, "Начну с проблематики, которую решает работа. Существует много различных метрик оценки качества изображений. Многие  из них устроены как последовательность осмысленных этапов. Однако  наиболее распространнённые реализации этих метрик, например в OPEN_CV, возвращают, лишь итоговое значение. Не давая возможности посмотреть  промежуточные вычисления. Из-за этого трудно понять, какие именно  свойства фрагмента повлияли на оценку алгоритма. Поэтому **цель работы** разработка desktop-приложения для визуализации алгоритмов IQA. Пользователь загружает снимок, выделяет область интереса и получает **поэтапное представление расчёта**. с модульной архитектурой, для возможности интегрировать новые алгоритмы. Чтобы достичь этой цели, были решины  задачи — от изучнеия теории и выбора стека до реализации ядра, интерфейса и проверки на тестовых изображениях.")

# ===========================================================================
# Слайд 3 — Оценка качества изображений (IQA)
# ===========================================================================
s = slide()
header(s, "Оценка качества изображений (IQA)", kicker="ТЕОРИЯ", number="3")
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.6))
para(tf, "IQA — численная оценка качества, согласованная с восприятием человека (MOS). "
         "По наличию эталона методы делят на три класса:",
     14, color=MUTE, space_after=0)
iqa = [
    ("Full-Reference (FR)", "нужен оригинал", "сравнение кодеков, эталонные тесты"),
    ("Reduced-Reference (RR)", "часть признаков оригинала", "передача по ограниченному каналу"),
    ("No-Reference (NR)", "эталон не нужен", "видеонаблюдение, архивы, реальное время"),
]
tx = Inches(0.55)
ty = Inches(2.05)
c1, c2, c3 = Inches(3.7), Inches(3.3), Inches(5.2)
cx = tx
for h, w in [("Подход", c1), ("Эталон", c2), ("Область применения", c3)]:
    add_rect(s, cx, ty, w, Inches(0.55), NAVY)
    tb, tf = textbox(s, cx + Inches(0.18), ty, w - Inches(0.3), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, h, 14, color=WHITE, bold=True, space_after=0)
    cx += w
ty += Inches(0.55)
rh = Inches(0.95)
for i, (a, b, c) in enumerate(iqa):
    is_nr = (i == 2)
    bg = RGBColor(0xE8, 0xF1, 0xFA) if is_nr else (WHITE if i % 2 == 0 else LIGHT)
    cx = tx
    for val, w in [(a, c1), (b, c2), (c, c3)]:
        add_rect(s, cx, ty, w, rh, bg)
        tb, tf = textbox(s, cx + Inches(0.18), ty, w - Inches(0.32), rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, val, 14, color=(ACCENT if (val == a and is_nr) else INK),
             bold=(val == a), space_after=0)
        cx += w
    ty += rh
add_rect(s, tx, ty + Inches(0.15), c1 + c2 + c3, Inches(0.72), NAVY)
add_rect(s, tx, ty + Inches(0.15), Inches(0.1), Inches(0.72), ACCENT)
tb, tf = textbox(s, tx + Inches(0.3), ty + Inches(0.15), Inches(11.6), Inches(0.72), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "NR-метрики работают без оригинала " + ARR + " подходят для реальных задач. "
         "BRISQUE использует статистику естественных сцен (NSS) в пространственной области.",
     14, color=WHITE, bold=True, space_after=0)
footer(s)
notes(s, "Несколько слов о предметной области. Задача IQA — численно оценить качество "
         "изображения так, чтобы оценка согласовывалась с субъективным восприятием человека, "
         "то есть со средними экспертными баллами MOS. По наличию эталона методы делятся на три "
         "класса. Full-Reference требует исходное неискажённое изображение и применяется для "
         "сравнения кодеков. Reduced-Reference использует лишь часть признаков оригинала. "
         "No-Reference вообще не требует эталона — именно этот класс нужен в реальных задачах: "
         "видеонаблюдение, обработка архивов, работа в реальном времени, где оригинала нет. "
         "Поэтому в работе рассматривается NR-подход, а в качестве примера — BRISQUE, основанный "
         "на статистике естественных сцен.")

# ===========================================================================
# Слайд 4 — Технологический стек
# ===========================================================================
s = slide()
header(s, "Технологический стек", kicker="ИНСТРУМЕНТЫ", number="4")
stack = [
    ("TypeScript", "строгая типизация данных пайплайна между worker и UI"),
    ("Electron", "локальный desktop-доступ к файлам без сервера"),
    ("Web Worker", "тяжёлые свёртки и признаки — без блокировки интерфейса"),
    ("HTML Canvas + DOM", "карты, гистограммы и таблицы без тяжёлых библиотек"),
    ("electron-vite", "сборка renderer и worker"),
]
y = Inches(1.6)
for name, desc in stack:
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), LIGHT)
    add_rect(s, Inches(0.55), y, Inches(3.25), Inches(0.92), NAVY)
    tb, tf = textbox(s, Inches(0.7), y, Inches(3.0), Inches(0.92), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, 18, color=WHITE, bold=True, space_after=0)
    tb, tf = textbox(s, Inches(4.1), y, Inches(8.4), Inches(0.92), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, desc, 16, color=INK, space_after=0)
    y += Inches(1.04)
footer(s)
notes(s, "Стек подобран под интерактивный сценарий. TypeScript снижает ошибки при передаче "
         "массивов между worker и интерфейсом. Electron даёт работу с локальными файлами. "
         "Тяжёлые свёртки и извлечение признаков вынесены в Web Worker — интерфейс остаётся "
         "отзывчивым при перетаскивании ROI. Визуализация — на Canvas и в DOM, без тяжёлых "
         "графических библиотек.")

# ===========================================================================
# Слайд 5 — Разметка и запуск (ui-main.png)
# ===========================================================================
s = slide()
header(s, "Разметка и запуск приложения", kicker="ТОЧКА ВХОДА", number="5")
bx, btf = textbox(s, Inches(0.55), Inches(1.55), Inches(5.1), Inches(5.2))
b_items = [
    "index.html — статичная оболочка окна",
    "Пустые контейнеры #qa-tabs-nav, #qa-tabs-container",
    "renderer.ts — точка входа, requireElement (fail-fast)",
    "Сборка ShellElements " + ARR + " new AppController(els)",
    "Вкладки метода не зашиты в HTML",
]
for it in b_items:
    para(btf, it, 17, bullet=True, space_after=14)
img = os.path.join(FIG, "ui-main.png")
if os.path.exists(img):
    iw = Inches(6.85)
    ih = iw * 983 / 1482
    ix = Inches(5.95)
    iy = Inches(1.95)
    add_rect(s, ix - Inches(0.06), iy - Inches(0.06), iw + Inches(0.12), ih + Inches(0.12), LINE)
    s.shapes.add_picture(img, ix, iy, width=iw, height=ih)
footer(s)
notes(s, "Интерфейс начинается с разметки index.html: workspace для снимка, сайдбар, canvas "
         "превью, блок оценки. Важно — вкладки «Карты», «Графики», «Признаки» в HTML не "
         "прописаны, только пустые контейнеры. Точка входа renderer.ts находит все нужные "
         "элементы через requireElement и передаёт их в AppController. Если критичный узел "
         "отсутствует, приложение падает сразу — так проще отлавливать ошибки разметки. "
         "AppController — центральный координатор всего сценария работы пользователя.")

# ===========================================================================
# Слайд 6 — Классы UI в AppController (таблица)
# ===========================================================================
s = slide()
header(s, "Классы UI в AppController", kicker="КООРДИНАЦИЯ ИНТЕРФЕЙСА", number="6")
rows = [
    ("ViewportManager", "zoom, scroll, загрузка снимка"),
    ("SelectionManager", "рисование и изменение ROI"),
    ("ScorePresenter", "итоговая оценка в шапке сайдбара"),
    ("TabHost", "динамические вкладки сайдбара"),
    ("FullscreenView", "увеличенный просмотр карт"),
    ("CanvasContextMenu", "экспорт canvas в PNG"),
    ("HelpPanel", "академическая справка"),
]
ty = Inches(1.5)
rh = Inches(0.62)
tx = Inches(0.55)
w1, w2 = Inches(4.1), Inches(8.1)
add_rect(s, tx, ty, w1, rh, NAVY)
add_rect(s, tx + w1, ty, w2, rh, NAVY)
tb, tf = textbox(s, tx + Inches(0.2), ty, w1 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Класс", 15, color=WHITE, bold=True, space_after=0)
tb, tf = textbox(s, tx + w1 + Inches(0.2), ty, w2 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Роль", 15, color=WHITE, bold=True, space_after=0)
ty += rh
for i, (cls, role) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, tx, ty, w1, rh, bg)
    add_rect(s, tx + w1, ty, w2, rh, bg)
    tb, tf = textbox(s, tx + Inches(0.2), ty, w1 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, cls, 15, color=NAVY, bold=True, font="Consolas", space_after=0)
    tb, tf = textbox(s, tx + w1 + Inches(0.2), ty, w2 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, role, 15, color=INK, space_after=0)
    ty += rh
add_rect(s, tx, ty + Inches(0.12), w1 + w2, Inches(0.62), RGBColor(0xE8, 0xF1, 0xFA))
tb, tf = textbox(s, tx + Inches(0.25), ty + Inches(0.12), w1 + w2 - Inches(0.5),
                 Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Изменение ROI " + ARR + " превью на previewCanvas " + ARR + " запрос в worker. "
         "AppController оркестрирует UI, не зная деталей BRISQUE.",
     14, color=NAVY, bold=True, space_after=0)
footer(s)
notes(s, "В конструкторе AppController создаются сервисы отображения. ViewportManager отвечает "
         "за масштаб и прокрутку workspace. SelectionManager — за интерактивную рамку ROI. "
         "Когда пользователь меняет выделение, фрагмент копируется на preview-canvas, и "
         "AppController ставит задачу анализа. ScorePresenter показывает итоговый score. "
         "TabHost управляет вкладками — до первого анализа они скрыты. FullscreenView и "
         "CanvasContextMenu — общие для любого метода. AppController не знает деталей BRISQUE — "
         "он только оркестрирует UI и обмен с worker.")

# ===========================================================================
# Слайд 7 — Модульная оболочка и контракты
# ===========================================================================
s = slide()
header(s, "Модульная оболочка и контракты", kicker="АРХИТЕКТУРА (1/2)", number="7")
contracts = [
    ("QualityMethod", "создаёт worker и разбирает ответ"),
    ("MethodUiPlugin", "вкладки сайдбара, fullscreen, export"),
    ("MethodRegistry", "реестр методов (BRISQUE, placeholder)"),
    ("MethodSwitcher", "смена метода: dispose() и новый worker"),
    ("AnalysisResult", "единый формат: score + payload"),
]
y = Inches(1.55)
for name, desc in contracts:
    add_rect(s, Inches(0.55), y, Inches(8.0), Inches(0.82), LIGHT)
    add_rect(s, Inches(0.55), y, Inches(0.08), Inches(0.82), ACCENT)
    tb, tf = textbox(s, Inches(0.8), y, Inches(3.1), Inches(0.82), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, 16, color=NAVY, bold=True, font="Consolas", space_after=0)
    tb, tf = textbox(s, Inches(3.95), y, Inches(4.4), Inches(0.82), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, desc, 14, color=INK, space_after=0)
    y += Inches(0.93)
add_rect(s, Inches(8.85), Inches(1.55), Inches(3.9), Inches(4.55), NAVY)
tb, tf = textbox(s, Inches(9.15), Inches(1.85), Inches(3.35), Inches(4.0))
para(tf, "ЗАДЕЛ", 13, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True, space_after=10)
para(tf, "Новая метрика =\nновая пара\nmethod + plugin", 19, color=WHITE, bold=True,
     space_after=14, line_spacing=1.1)
para(tf, "index.html менять не нужно", 15, color=RGBColor(0xCF, 0xDD, 0xEC), space_after=12)
para(tf, "Уже подключены:\nBRISQUE и placeholder", 14, color=RGBColor(0xCF, 0xDD, 0xEC),
     space_after=0, line_spacing=1.1)
footer(s)
notes(s, "Чтобы BRISQUE не «прирос» к интерфейсу, введена модульная оболочка в каталоге shell. "
         "Вычислительная часть реализует QualityMethod: создаёт worker и превращает ответ в "
         "AnalysisResult. Визуальная — MethodUiPlugin: монтирует вкладки и подключает fullscreen. "
         "MethodRegistry хранит зарегистрированные методы. При смене алгоритма MethodSwitcher "
         "вызывает dispose у старого плагина, очищает вкладки и поднимает новый worker. Оболочка "
         "всегда работает с одним форматом результата — score и произвольный payload.")

# ===========================================================================
# Слайд 8 — Схема потока данных
# ===========================================================================
LX = Inches(1.25)       # левый край слоя
LW = Inches(10.83)      # ширина слоя
AX = int((EMUW - Inches(0.42)) / 2)  # центр вертикальной стрелки
ARROW_BLUE = RGBColor(0x7E, 0x9B, 0xB8)
BOX_BLUE = RGBColor(0xE8, 0xF1, 0xFA)


def container(x, y, w, h, bg, border, title, title_color, dashed=False):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg
    sp.line.color.rgb = border
    sp.line.width = Pt(1.25)
    if dashed:
        try:
            from pptx.enum.line import MSO_LINE_DASH_STYLE
            sp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    sp.shadow.inherit = False
    tb, tf = textbox(s, x + Inches(0.22), y + Inches(0.05), w - Inches(0.44), Inches(0.32),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 12, color=title_color, bold=True, space_after=0)


def minibox(x, y, w, h, lines, bg, fg, size=11, font=FONT, border=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg
    if border is not None:
        sp.line.color.rgb = border
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    tb, tf = textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    for j, t in enumerate(lines):
        para(tf, t, size if j == 0 else size - 1, color=fg, bold=(j == 0),
             align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.95,
             font=font if j == 0 else FONT)


def rarrow(xc, yc):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(xc - Inches(0.11)),
                           int(yc - Inches(0.08)), Inches(0.22), Inches(0.16))
    fill(a, ARROW_BLUE)
    a.shadow.inherit = False


def varrow(ytop, h, label1=None, label2=None):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, AX, ytop, Inches(0.42), h)
    fill(a, ARROW_BLUE)
    a.shadow.inherit = False
    if label1:
        tb, tf = textbox(s, Inches(7.35), int(ytop - Inches(0.04)), Inches(4.7),
                         int(h + Inches(0.08)), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, label1, 11, color=NAVY, bold=True, space_after=1, font="Consolas")
        if label2:
            para(tf, label2, 10, color=MUTE, italic=True, space_after=0)


s = slide()
header(s, "Схема потока данных", kicker="АРХИТЕКТУРА (2/2)", number="8")

# --- Слой 1: главный поток renderer (запрос) ---
container(LX, Inches(1.30), LW, Inches(1.12), BOX_BLUE, NAVY,
          "ГЛАВНЫЙ ПОТОК · renderer (AppController)", NAVY)
mb_y = Inches(1.70)
mb_h = Inches(0.62)
a_boxes = [
    (Inches(1.50), ["SelectionManager", "выделение ROI (crop)"]),
    (Inches(5.15), ["previewCanvas", "getImageData " + ARR + " RGBA"]),
    (Inches(8.80), ["postMessage", "(RGBA, requestId)"]),
]
for bx, lines in a_boxes:
    minibox(bx, mb_y, Inches(3.1), mb_h, lines, WHITE, NAVY, size=12, font="Consolas",
            border=RGBColor(0xBF, 0xD4, 0xE8))
rarrow(Inches(4.875), int(mb_y + mb_h / 2))
rarrow(Inches(8.525), int(mb_y + mb_h / 2))

varrow(Inches(2.46), Inches(0.32),
       "postMessage(RGBA, requestId)", "Transferable: rgbaArray.buffer")

# --- Слой 2: Web Worker (отдельный поток) ---
container(LX, Inches(2.82), LW, Inches(1.86),
          RGBColor(0xF4, 0xF4, 0xF4), MUTE,
          "WEB WORKER · brisque.worker  " + ARR + "  BrisquePipeline.execute  (отдельный поток)",
          RGBColor(0x44, 0x4C, 0x55), dashed=True)
pipe = [
    (["Grayscale", "BT.601"], NAVY),
    (["MSCN", "Scale 1 / 2"], NAVY),
    (["Pairwise", "H·V·D" + SUB1 + "·D" + SUB2], NAVY),
    (["Признаки ×36", "GGD / AGGD"], NAVY),
    (["SVR"], NAVY),
    (["DMOS", "(score)"], ACCENT),
]
pb_y = Inches(3.45)
pb_h = Inches(0.92)
pb_w = Inches(1.45)
pb_gap = Inches(0.34)
px = Inches(1.5)
for k, (lines, bg) in enumerate(pipe):
    minibox(px, pb_y, pb_w, pb_h, lines, bg, WHITE, size=11)
    if k < len(pipe) - 1:
        rarrow(int(px + pb_w + pb_gap / 2), int(pb_y + pb_h / 2))
    px = int(px + pb_w + pb_gap)

varrow(Inches(4.72), Inches(0.32),
       "ответ: BrisqueWorkerSuccess", "Transferable: карты, признаки, score")

# --- Слой 3: главный поток renderer (разбор и вывод) ---
container(LX, Inches(5.08), LW, Inches(1.44), BOX_BLUE, NAVY,
          "ГЛАВНЫЙ ПОТОК · renderer (AppController)", NAVY)
minibox(Inches(1.5), Inches(5.44), Inches(10.33), Inches(0.46),
        ["parseWorkerMessage  " + ARR + "  AnalysisResult  (score + payload)"],
        NAVY, WHITE, size=12, font="Consolas")
minibox(Inches(1.5), Inches(5.98), Inches(4.6), Inches(0.46),
        ["ScorePresenter  " + ARR + "  score в шапке"],
        WHITE, NAVY, size=11, font="Consolas", border=RGBColor(0xBF, 0xD4, 0xE8))
minibox(Inches(6.4), Inches(5.98), Inches(5.43), Inches(0.46),
        ["TabHost.dispatchResult  " + ARR + "  «Карты» · «Графики» · «Признаки»"],
        WHITE, NAVY, size=11, font="Consolas", border=RGBColor(0xBF, 0xD4, 0xE8))
footer(s)
notes(s, "Эта схема показывает, что куда и откуда передаётся; пунктиром выделена граница потоков. "
         "Сверху — главный поток renderer. Когда пользователь меняет выделение, SelectionManager "
         "отдаёт прямоугольник ROI, AppController рисует фрагмент на previewCanvas, считывает "
         "пиксели через getImageData и отправляет их в воркер вызовом postMessage. Массив RGBA "
         "передаётся как Transferable — без копирования; для защиты от перегрузки используется "
         "requestId. В середине — Web Worker, отдельный поток: тонкий адаптер brisque.worker "
         "запускает BrisquePipeline.execute, который считает яркость по BT.601, карты MSCN и "
         "попарные произведения на двух масштабах, извлекает 36 признаков через GGD и AGGD и "
         "прогоняет их через SVR, получая итоговую оценку DMOS. Результат — объект "
         "BrisqueWorkerSuccess с картами, признаками и score — возвращается обратно, тоже как "
         "Transferable. Снова в главном потоке BrisqueQualityMethod.parseWorkerMessage превращает "
         "его в единый AnalysisResult со score и payload. Дальше ScorePresenter показывает оценку "
         "в шапке, а TabHost.dispatchResult раздаёт payload в панели «Карты», «Графики» и "
         "«Признаки».")

# ===========================================================================
# Слайд 9 — BRISQUE: этапы расчёта (цепочка)
# ===========================================================================
s = slide()
header(s, "BRISQUE: этапы расчёта", kicker="АЛГОРИТМ (1/2)", number="9")
chain = [
    "Яркость\n(BT.601)",
    "MSCN\n(" + MU + ", " + SIG + ", " + SIGH + ")",
    "Pairwise\nH, V, D" + SUB1 + ", D" + SUB2,
    "GGD / AGGD",
    "36 признаков\n(2 масштаба)",
    "SVR",
    "score",
]
n = len(chain)
total_w = Inches(12.2)
box_w = Inches(1.5)
arrow_w = (total_w - box_w * n) / (n - 1)
x = Inches(0.55)
y = Inches(2.35)
box_h = Inches(1.3)
for i, label in enumerate(chain):
    is_last = (i == n - 1)
    color = ACCENT if is_last else NAVY
    add_rect(s, x, y, box_w, box_h, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, x, y, box_w, box_h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
         space_after=0, line_spacing=1.0)
    if i < n - 1:
        ax = x + box_w
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax + Emu(int(arrow_w * 0.15)),
                                 y + box_h / 2 - Inches(0.12),
                                 Emu(int(arrow_w * 0.7)), Inches(0.24))
        fill(arr, RGBColor(0x9A, 0xB4, 0xCC))
        arr.shadow.inherit = False
    x += box_w + arrow_w
mapping = [
    ("Карты", "вкладка «Карты»"),
    ("Гистограммы + кривые", "вкладка «Графики»"),
    ("Таблица " + SIG + ", " + NU + SUP2 + ", " + LAM, "вкладка «Признаки»"),
]
mw = Inches(3.95)
mx = Inches(0.55)
my = Inches(4.6)
for title, sub in mapping:
    add_rect(s, mx, my, mw, Inches(1.2), LIGHT)
    add_rect(s, mx, my, Inches(0.08), Inches(1.2), ACCENT)
    tb, tf = textbox(s, mx + Inches(0.25), my + Inches(0.18), mw - Inches(0.5), Inches(0.9))
    para(tf, title, 16, color=INK, bold=True, space_after=4)
    para(tf, sub, 13, color=MUTE, space_after=0)
    mx += mw + Inches(0.2)
footer(s)
notes(s, "BRISQUE работает в пространственной области. Сначала RGB переводится в яркость, затем "
         "локальная нормализация MSCN выделяет текстуру независимо от освещения. Попарные "
         "произведения соседних коэффициентов фиксируют пространственные связи. Для MSCN "
         "подбирается GGD, для pairwise — AGGD. На двух масштабах формируется 36 признаков, и "
         "SVR выдаёт итоговую оценку, обученную на субъективных данных LIVE. Каждый этап "
         "отображается в интерфейсе: карты на Canvas, распределения на графиках, числа — в "
         "таблице признаков.")

# ===========================================================================
# Слайд 10 — BRISQUE: математическая модель (формулы)
# ===========================================================================
F_MSCN = render_formula("mscn",
    r"$\hat{I}(x,y)=\frac{I(x,y)-\mu(x,y)}{\sigma(x,y)+C}$", fontsize=30)
F_GGD = render_formula("ggd",
    r"$f(x;\,\alpha,\sigma^2)=\frac{\alpha}{2\beta\,\Gamma(1/\alpha)}\,\mathrm{exp}\!\left[-\left(\frac{|x|}{\beta}\right)^{\alpha}\right]$",
    fontsize=30)
F_AGGD = render_formula("aggd",
    r"$f(x)=\frac{\alpha}{(\beta_l+\beta_r)\,\Gamma(1/\alpha)}\,\mathrm{exp}\!\left[-\left(\frac{|x|}{\beta_{l,r}}\right)^{\alpha}\right]$",
    fontsize=30)

s = slide()
header(s, "BRISQUE: математическая модель", kicker="АЛГОРИТМ (2/2)", number="10")
bands = [
    ("1. Нормализация локального контраста (MSCN)", F_MSCN,
     MU + ", " + SIG + " — локальные среднее и СКО по гауссову окну 7" + TIMES + "7;   C = 1 — стабилизация"),
    ("2. Аппроксимация гистограммы MSCN — GGD", F_GGD,
     "\u03b2 = " + SIG + RAD + "(" + GAM + "(1/" + ALP + ")/" + GAM + "(3/" + ALP + "));   признаки (" + ALP + ", " + SIG + SUP2 + ")"),
    ("3. Попарные произведения H, V, D" + SUB1 + ", D" + SUB2 + " — AGGD", F_AGGD,
     "вектор (" + ALP + ", " + ETA + ", " + SIG + "\u2097" + SUP2 + ", " + SIG + "\u1d63" + SUP2 + ") " + ARR +
     " 18 признаков " + TIMES + " 2 масштаба = 36"),
]
by = Inches(1.42)
bh = Inches(1.78)
gap = Inches(0.12)
for label, finfo, sub in bands:
    add_rect(s, Inches(0.55), by, Inches(12.2), bh, LIGHT)
    add_rect(s, Inches(0.55), by, Inches(0.1), bh, ACCENT)
    tb, tf = textbox(s, Inches(0.85), by + Inches(0.1), Inches(11.6), Inches(0.4))
    para(tf, label, 15, color=NAVY, bold=True, space_after=0)
    # формула по центру полосы
    place_formula(s, finfo, Inches(0.85), by + Inches(0.45), Inches(11.6),
                  Inches(0.95), max_h_in=0.62)
    tb, tf = textbox(s, Inches(0.85), by + bh - Inches(0.36), Inches(11.6), Inches(0.32))
    para(tf, sub, 11.5, color=MUTE, align=PP_ALIGN.CENTER, space_after=0)
    by += bh + gap
footer(s)
notes(s, "Поскольку работа математической направленности, кратко о модели. Карта MSCN — это "
         "нормализация локального контраста: из яркости вычитается локальное среднее mu и "
         "результат делится на локальное СКО sigma плюс константа. Для естественных снимков "
         "гистограмма MSCN близка к гауссовой, а искажения её меняют. Эту форму описывает "
         "обобщённое гауссово распределение GGD с параметром формы alpha и дисперсией. Попарные "
         "произведения соседних коэффициентов асимметричны, поэтому для них применяется "
         "асимметричное распределение AGGD с четырьмя параметрами на каждое из направлений. "
         "Итого 18 признаков на масштаб, на двух масштабах — 36, которые подаются в SVR.")

# ===========================================================================
# Слайд 11 — Демонстрация (крупные скриншоты)
# ===========================================================================
s = slide()
header(s, "Демонстрация: вкладки в действии", kicker="РЕЗУЛЬТАТЫ (1/2)", number="11")
shots = [
    ("viz-map-mscn.png", "Карта MSCN"),
    ("viz-chart-mscn.png", "Гистограмма + кривая GGD"),
    ("viz-features.png", "Таблица 36 признаков"),
]
iw = Inches(4.18)
ih = iw * 983 / 1482
gapx = Inches(0.1)
block_w = iw * 3 + gapx * 2
ix0 = int((EMUW - block_w) / 2)
iy = Inches(2.05)
ix = ix0
for fn, cap in shots:
    p = os.path.join(FIG, fn)
    if os.path.exists(p):
        add_rect(s, ix - Inches(0.05), iy - Inches(0.05), iw + Inches(0.1), ih + Inches(0.1), LINE)
        s.shapes.add_picture(p, int(ix), int(iy), width=int(iw), height=int(ih))
    tb, tf = textbox(s, int(ix), int(iy + ih + Inches(0.08)), int(iw), Inches(0.4))
    para(tf, cap, 14, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    ix += iw + gapx
footer(s)
notes(s, "Покажу инструмент в действии на трёх вкладках. Слева — карта коэффициентов MSCN для "
         "выделенного фрагмента: видно текстуру независимо от освещения. В центре — гистограмма "
         "значений MSCN с наложенной подогнанной кривой GGD; по ней наглядна форма распределения. "
         "Справа — таблица всех 36 числовых признаков, попадающих в регрессор. Любую карту можно "
         "открыть на весь экран и сохранить в PNG.")

# ===========================================================================
# Слайд 12 — Тестирование (серия из 25 снимков, сравнение с MATLAB)
# ===========================================================================
def render_delta_chart():
    deltas = [1.94, 1.90, 3.77, 3.51, 3.68, 1.42, 4.24, 2.66, 3.53, 2.05,
              2.20, 2.21, 1.64, 3.24, 2.02, 1.96, 4.42, 2.51, 1.95, 0.85,
              3.15, 1.94, 1.42, 2.08, 1.73]
    path = os.path.join(FORM_DIR, "delta_chart.png")
    fig, ax = plt.subplots(figsize=(6.4, 3.85), dpi=200)
    fig.patch.set_alpha(0)
    ax.patch.set_alpha(0)
    xs = list(range(1, 26))
    ax.bar(xs, deltas, width=0.72, color="#1E3A5F")
    ax.axhline(2.51, color="#2E86C1", linestyle="--", linewidth=1.8)
    ax.text(25.4, 2.66, "среднее 2.51", color="#2E86C1", fontsize=10,
            ha="right", va="bottom")
    ax.set_xticks([1, 5, 10, 15, 20, 25])
    ax.set_xlabel("\u2116 снимка (I01\u2013I25)", fontsize=10)
    ax.set_ylabel("|\u0394|, балла", fontsize=10)
    ax.set_xlim(0.3, 25.7)
    ax.set_ylim(0, 5)
    ax.tick_params(labelsize=9)
    for sp in ("top", "right"):
        ax.spines[sp].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, dpi=200, transparent=True)
    plt.close(fig)
    w, h = Image.open(path).size
    return path, w, h


def statcard(x, y, w, h, big, small, big_size=22):
    add_rect(s, x, y, w, h, LIGHT)
    add_rect(s, x, y, Inches(0.08), h, ACCENT)
    tb, tf = textbox(s, x + Inches(0.22), y, Inches(2.05), h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, big, big_size, color=NAVY, bold=True, space_after=0)
    tb, tf = textbox(s, x + Inches(2.35), y, w - Inches(2.5), h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, small, 12.5, color=INK, space_after=0, line_spacing=1.0)


s = slide()
header(s, "Тестирование: серия из 25 снимков vs MATLAB", kicker="РЕЗУЛЬТАТЫ (2/2)", number="12")

# --- столбчатая диаграмма |Δ| по 25 снимкам ---
cpath, cwpx, chpx = render_delta_chart()
cx0, cy0, cw = Inches(0.55), Inches(1.45), Inches(6.4)
chh = int(int(cw) * chpx / cwpx)
add_rect(s, cx0 - Inches(0.05), cy0 - Inches(0.05), cw + Inches(0.1), chh + Inches(0.1), LINE)
add_rect(s, cx0, cy0, cw, chh, WHITE)
s.shapes.add_picture(cpath, int(cx0), int(cy0), width=int(cw), height=chh)
tb, tf = textbox(s, int(cx0), int(cy0) + chh + Inches(0.06), int(cw), Inches(0.34),
                 anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Абсолютное расхождение |\u0394| с MATLAB по 25 полным кадрам",
     12, color=MUTE, align=PP_ALIGN.CENTER, space_after=0)

# --- карточки со статистикой ---
scx, scw, sch, scgap = Inches(7.15), Inches(5.65), Inches(0.66), Inches(0.12)
statcard(scx, Inches(1.45), scw, sch, "25", "тестовых снимков \u00b7 полный кадр")
statcard(scx, Inches(1.45) + (sch + scgap), scw, sch, "2.51", "среднее |\u0394| (балла)")
statcard(scx, Inches(1.45) + (sch + scgap) * 2, scw, sch, "0.85\u20134.42",
         "диапазон |\u0394|: мин I20, макс I17", big_size=18)

# --- два иллюстративных ROI ---
roi = [
    ("test-good.png", "I14 \u2014 качественный", "9.09 vs 12.33  (|\u0394| 3.24)"),
    ("test-bad.png", "I20 \u2014 искажённый", "92.08 vs 89.35  (|\u0394| 2.73)"),
]
riw = Inches(2.6)
rih = int(int(riw) * 384 / 512)
rgap = Inches(0.25)
rxs = [scx, scx + riw + rgap]
ry = Inches(3.85)
for (fn, cap, sc), rx in zip(roi, rxs):
    p = os.path.join(FIG, fn)
    if os.path.exists(p):
        add_rect(s, rx - Inches(0.04), ry - Inches(0.04), riw + Inches(0.08), rih + Inches(0.08), LINE)
        s.shapes.add_picture(p, int(rx), int(ry), width=int(riw), height=rih)
    tb, tf = textbox(s, int(rx), int(ry) + rih + Inches(0.04), int(riw), Inches(0.6))
    para(tf, cap, 12, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=1)
    para(tf, sc, 11, color=MUTE, align=PP_ALIGN.CENTER, space_after=0)

# --- итоговая лента ---
vy = Inches(6.18)
add_rect(s, Inches(0.55), vy, Inches(12.23), Inches(0.62), RGBColor(0xE8, 0xF1, 0xFA))
tb, tf = textbox(s, Inches(0.75), vy, Inches(11.83), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Ранжирование верное: I20 (92.08) " + GG + " I14 (9.09).  По 25 кадрам среднее "
         "|\u0394| = 2.51 \u2014 корректность вычислительного ядра подтверждена",
     13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s)
notes(s, "Корректность ядра проверена на серии из 25 изображений I01\u2013I25: каждый полный кадр "
         "анализировался в приложении и в эталонной реализации BRISQUE на MATLAB по идентичным "
         "пикселям. На диаграмме слева видно абсолютное расхождение score по каждому снимку: все "
         "значения укладываются в небольшой коридор, среднее отклонение \u2014 2.51 балла, максимум "
         "4.42 (I17), минимум 0.85 (I20). Справа \u2014 ключевые цифры и два характерных ROI: на "
         "качественном фрагменте оценка около девяти, на искажённом \u2014 около девяноста; "
         "ранжирование полностью совпадает с визуальным качеством. Остаточные расхождения "
         "объясняются отличиями бикубического масштабирования, граничной обработки и арифметики с "
         "плавающей точкой (float32).")

# ===========================================================================
# Слайд 13 — Итоги
# ===========================================================================
s = slide()
add_rect(s, 0, 0, EMUW, EMUH, NAVY)
add_rect(s, 0, 0, Inches(0.22), EMUH, ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9))
para(tf, "ИТОГИ", 14, color=RGBColor(0xAE, 0xC8, 0xE0), bold=True, space_after=4)
para(tf, "Результаты и перспективы", 30, color=WHITE, bold=True, space_after=0)
add_rect(s, Inches(0.92), Inches(1.95), Inches(2.6), Inches(0.05), ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(3.0))
para(tf, "Цель достигнута: инструмент для пошаговой визуализации IQA на выбранном ROI",
     20, color=WHITE, bullet=True, space_after=16)
para(tf, "Модульная архитектура — задел на NIQE, PIQE и другие метрики без переработки оболочки",
     20, color=WHITE, bullet=True, space_after=16)
para(tf, "Перспективы: пакетная обработка, расширенный экспорт результатов",
     20, color=WHITE, bullet=True, space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2))
para(tf, "Спасибо за внимание!", 30, color=WHITE, bold=True, space_after=4)
para(tf, "Готов ответить на ваши вопросы", 17, color=RGBColor(0xCF, 0xDD, 0xEC), space_after=0)
notes(s, "В результате разработано desktop-приложение, которое позволяет интерактивно "
         "исследовать безэталонную метрику на выбранном фрагменте изображения. Реализованы "
         "вычислительное ядро, модульный интерфейс и проверка на тестовых данных. Заложенная "
         "архитектура допускает подключение других метрик без переработки оболочки. Спасибо за "
         "внимание, готов ответить на вопросы.")

# ===========================================================================
# Слайд 14 — Литература
# ===========================================================================
s = slide()
header(s, "Список литературы", kicker="ИСТОЧНИКИ", number="14")
refs = [
    "Zhai G., Min X. Perceptual image quality assessment: a survey. Science China, 2020.",
    "Mittal A., Moorthy A. K., Bovik A. C. No-Reference IQA in the Spatial Domain (BRISQUE). IEEE TIP, 2012.",
    "Sheikh H. R. et al. A Statistical Evaluation of Recent FR IQA Algorithms (LIVE). IEEE TIP, 2006.",
    "Saad M. A., Bovik A. C., Charrier C. Blind IQA: NSS Approach in the DCT Domain. IEEE TIP, 2012.",
    "Рекомендация ITU-R BT.601-6. Studio encoding parameters of digital television.",
    "Srivastava A. et al. On Advances in Statistical Modeling of Natural Images. JMIV, 2003.",
    "Ruderman D. L. The statistics of natural images. Network: Comp. in Neural Systems, 1994.",
    "Sharifi K., Leon-Garcia A. Estimation of Shape Parameter for GGD. IEEE TCSVT, 1995.",
    "Keys R. Cubic Convolution Interpolation for Digital Image Processing. IEEE TASSP, 1981.",
    "Chang C.-C., Lin C.-J. LIBSVM: A Library for Support Vector Machines. ACM TIST, 2011.",
    "Electron Documentation; WHATWG HTML Living Standard — Web Workers, Canvas.",
    "TypeScript Handbook; MathWorks Documentation — imresize.",
]
half = (len(refs) + 1) // 2
columns = [refs[:half], refs[half:]]
col_x = [Inches(0.55), Inches(6.9)]
for ci, items in enumerate(columns):
    tb, tf = textbox(s, col_x[ci], Inches(1.55), Inches(5.9), Inches(5.3))
    for k, ref in enumerate(items):
        idx = ci * half + k + 1
        para(tf, str(idx) + ".  " + ref, 12, color=INK, space_after=10, line_spacing=1.04)
footer(s)
notes(s, "Ключевые источники: оригинальная статья BRISQUE (Mittal и др., 2012), база LIVE для "
         "обучения регрессора, работы по статистике естественных изображений (Ruderman, "
         "Srivastava), методы оценки параметров GGD (Sharifi), LIBSVM для SVR, а также "
         "техническая документация Electron, Web Workers, Canvas и TypeScript.")

prs.save(OUT)
print("OK ->", OUT)
print("slides:", len(prs.slides._sldIdLst))
